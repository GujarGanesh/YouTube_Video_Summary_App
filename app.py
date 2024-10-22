import streamlit as st
from dotenv import load_dotenv
import os
import re
import google.generativeai as genai
from youtube_transcript_api import YouTubeTranscriptApi
from gtts import gTTS
import tempfile
from PyPDF2 import PdfReader

# Load environment variables
load_dotenv()

# Configure Google Gemini API
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Function to generate the summarization prompt based on the selected style
def get_summarization_prompt(style):
    if style == "Bullet Points":
        return "Summarize the video in clear bullet points (max 250 words): "
    elif style == "Detailed Paragraph":
        return "Provide a detailed paragraph summarizing the video (max 250 words): "
    elif style == "Short Summary":
        return "Summarize the video briefly (max 100 words): "
    elif style == "Key Highlights":
        return "Highlight key points and action items from the video (max 150 words): "

# Function to extract the YouTube video ID from various URL formats
def extract_video_id(youtube_video_url):
    pattern = r"(?:v=|\/)([0-9A-Za-z_-]{11}).*"
    match = re.search(pattern, youtube_video_url)
    if match:
        return match.group(1)
    else:
        raise ValueError("Invalid YouTube URL format")

# Function to get the transcript data from YouTube videos
def extract_transcript_details(youtube_video_url):
    try:
        video_id = extract_video_id(youtube_video_url)
        transcript_text = YouTubeTranscriptApi.get_transcript(video_id)
        transcript = " ".join([i["text"] for i in transcript_text])
        return transcript
    except Exception as e:
        raise e

# Function to generate the summary using Google Gemini Pro
def generate_gemini_content(transcript_text, prompt):
    model = genai.GenerativeModel("gemini-pro")
    response = model.generate_content(prompt + transcript_text)
    return response.text

# Function to convert text to speech using Google Text-to-Speech (gTTS)
def text_to_speech(text):
    clean_text = re.sub(r"[^\w\s,.!?']", "", text)  # Clean the text from special characters
    tts = gTTS(text=clean_text, lang='en')  # Convert text to speech
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')  # Temporary audio file
    tts.save(temp_file.name)  # Save audio
    return temp_file.name  # Return file path


# Set page configuration
st.set_page_config(page_title="AI App Suite", layout="wide")

# Sidebar navigation with all pages
st.sidebar.image("Youtube_Logo.png", width=150)
st.sidebar.title("AI App Suite")
nav_option = st.sidebar.radio(
    "Navigation",
    [
        "📄 About Us",
        "🎥 Summarize YouTube Videos",
        "📑 PDF Summarizer",
        "📜 Chat to PDF",
        "📊 PPT Maker"
    ]
)

# About Us Page
if nav_option == "📄 About Us":
    st.title("Welcome to AI App Suite!")
    st.markdown("""
        ### Explore our AI-driven solutions:
        - **🎥 Summarize YouTube Videos**: Automatically generate concise summaries from YouTube videos.
        - **📑 PDF Summarizer**: Instantly convert and summarize PDF documents.
        - **📜 Chat to PDF**: Enter text, and instantly convert it into a shareable PDF.
        - **📊 PPT Maker**: Quickly generate professional PowerPoint presentations from your notes.

        Leverage the power of AI to save time and simplify your work. With an easy-to-use interface and cutting-edge tools, AI App Suite is your one-stop shop for productivity.
    """)

# Summarize YouTube Videos Page
elif nav_option == "🎥 Summarize YouTube Videos":
    st.title("YouTube Video Summary Tool")
    youtube_link = st.text_input("Enter YouTube Video Link:")

    summarization_style = st.selectbox(
        "Select Summarization Style:",
        ["Bullet Points", "Detailed Paragraph", "Short Summary", "Key Highlights"]
    )

    if youtube_link:
        try:
            video_id = extract_video_id(youtube_link)
            st.image(f"http://img.youtube.com/vi/{video_id}/0.jpg", use_column_width=True)
        except ValueError as e:
            st.error(str(e))

    if st.button("Get Detailed Notes"):
        try:
            transcript_text = extract_transcript_details(youtube_link)
            if transcript_text:
                summary_prompt = get_summarization_prompt(summarization_style)
                summary = generate_gemini_content(transcript_text, summary_prompt)
                st.markdown("## Detailed Notes:")
                st.write(summary)
                st.session_state.summary = summary
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")

    if 'summary' in st.session_state:
        st.markdown("### Summary Text:")
        st.write(st.session_state.summary)

        if st.button("Listen to Summary"):
            summary_text = st.session_state.summary
            audio_file_path = text_to_speech(summary_text)
            st.audio(audio_file_path, format='audio/mp3')
    st.title("How to Summarize YouTube Videos?")

    # Applying the CSS styles to the content
    st.markdown("""<div class="step-container">
            <h3>Follow These 3 Simple Steps to Summarize a YouTube Video:</h3>
            <p><strong>Step 1:</strong> Get the YouTube video link<br>
            Copy and paste the YouTube video link into NoteGPT.</p>
            <p><strong>Step 2:</strong> Generate a Summary<br>
            Click the "Generate Summary" button, and NoteGPT will fetch the transcript and summarize the YouTube video.</p>
            <p><strong>Step 3:</strong> Read the AI summary<br>
            Read the concise summary and save valuable time.</p>
        </div>""", unsafe_allow_html=True)
# PDF Summarizer Page
elif nav_option == "📑 PDF Summarizer":
    st.title("AI PDF Summarizer")
    st.markdown("""
    AI PDF Summarizer can summarize long PDF documents in seconds. 
    It can convert PDFs to text, allowing you to ask your PDF.
    The best PDF Summary tool helps you save time and learning faster and better.
    """)
    # File uploader for PDF files
    uploaded_file = st.file_uploader("Upload PDF", type="pdf")

    if uploaded_file:
        reader = PdfReader(uploaded_file)
        text = ""

        # Extract text from each page
        for page in reader.pages:
            text += page.extract_text()

        if text:
            st.text_area("Extracted Text", value=text, height=300)

            # Button to summarize the extracted text
            if st.button("Summarize PDF"):
                prompt = """You are an AI text summarizer. Please summarize the following text 
                            into important points within 250 words. Here is the text: """
                with st.spinner("Generating summary..."):
                    summary = generate_gemini_content(text, prompt)
                st.success("Summary Generated!")
                st.markdown("## Summary:")
                st.write(summary)
        else:
            st.warning("No text found in the PDF file.")

# Chat to PDF Page
elif nav_option == "📜 Chat to PDF":
    st.title("Chat with your any PDF")
    st.markdown("""
        Chat with any PDF using NoteGPT's ChatPDF app, along with its alternative solutions.
        You can ask your PDF documents questions, seek answers, summaries, and more.
        all free and accessible.
    """)
    chat_input = st.text_area("Enter your text to convert to PDF")

    if st.button("Generate PDF"):
        from fpdf import FPDF

        if chat_input:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, chat_input)
            pdf.output("chat_to_pdf.pdf")

            with open("chat_to_pdf.pdf", "rb") as file:
                st.download_button(label="Download PDF", data=file, file_name="chat_to_pdf.pdf", mime="application/pdf")
        else:
            st.warning("Please enter text to convert.")

# PPT Maker Page
elif nav_option == "📊 PPT Maker":
    st.title("PPT Maker")

    # Option to either input slide content manually or upload a file
    slide_option = st.radio("Choose how to add content to your slides:", ("Input Manually", "Upload Text File (.txt)"))

    # If user chooses to input slide content manually
    if slide_option == "Input Manually":
        st.subheader("Enter Slide Content:")

        # Get title and content for each slide
        slide_title = st.text_input("Enter Slide Title:")
        slide_content = st.text_area("Enter Slide Content:", height=200)

        # Button to generate the PPT
        if st.button("Generate PPT"):
            from pptx import Presentation

            # Create a presentation object
            prs = Presentation()

            # Add a slide with title and content layout
            slide_layout = prs.slide_layouts[1]  # 1 is for a title and content slide
            slide = prs.slides.add_slide(slide_layout)

            # Set slide title and content
            slide.shapes.title.text = slide_title
            slide.shapes.placeholders[1].text = slide_content

            # Save the presentation
            ppt_file = "generated_ppt.pptx"
            prs.save(ppt_file)

            # Provide download button
            with open(ppt_file, "rb") as file:
                st.download_button(
                    label="Download PPT",
                    data=file,
                    file_name="generated_ppt.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

    # If user chooses to upload a text file for content
    elif slide_option == "Upload Text File (.txt)":
        uploaded_file = st.file_uploader("Upload a Text File", type="txt")

        if uploaded_file:
            text = uploaded_file.read().decode("utf-8")

            # Split the text by lines and generate PPT
            if st.button("Generate PPT from File"):
                from pptx import Presentation

                # Create a presentation object
                prs = Presentation()

                # Add slides for each line of text
                for line in text.splitlines():
                    if line.strip():  # Only add non-empty lines as slides
                        slide_layout = prs.slide_layouts[1]  # Title and content layout
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = "Slide"
                        slide.shapes.placeholders[1].text = line

                # Save the presentation
                ppt_file = "generated_ppt_from_file.pptx"
                prs.save(ppt_file)

                # Provide download button
                with open(ppt_file, "rb") as file:
                    st.download_button(
                        label="Download PPT",
                        data=file,
                        file_name="generated_ppt_from_file.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

# Footer
st.sidebar.markdown("---")
st.sidebar.markdown("Copyright © 2024; Designed ❤️ by Code Rangers")
